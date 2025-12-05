"""
DataMigrate AI - Sales Presentation Generator
Generates PowerPoint presentations in 5 languages: EN, ES, PT, DA, DE

Requirements:
pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme
PRIMARY_BLUE = RGBColor(37, 99, 235)  # #2563EB
DARK_BLUE = RGBColor(30, 58, 138)     # #1E3A8A
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(243, 244, 246)
DARK_GRAY = RGBColor(55, 65, 81)

def add_title_slide(prs, title, subtitle, lang_data):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_BLUE
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
    sub_frame = sub_box.text_frame
    sub_para = sub_frame.paragraphs[0]
    sub_para.text = subtitle
    sub_para.font.size = Pt(24)
    sub_para.font.color.rgb = WHITE
    sub_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_BLUE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        para.text = f"• {bullet}"
        para.font.size = Pt(20)
        para.font.color.rgb = DARK_GRAY
        para.space_after = Pt(12)

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_BLUE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(0.5))
    left_frame = left_title_box.text_frame
    left_para = left_frame.paragraphs[0]
    left_para.text = left_title
    left_para.font.size = Pt(22)
    left_para.font.bold = True
    left_para.font.color.rgb = PRIMARY_BLUE

    # Left content
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4), Inches(4))
    left_content_frame = left_content.text_frame
    left_content_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            para = left_content_frame.paragraphs[0]
        else:
            para = left_content_frame.add_paragraph()
        para.text = f"• {item}"
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY
        para.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(0.5))
    right_frame = right_title_box.text_frame
    right_para = right_frame.paragraphs[0]
    right_para.text = right_title
    right_para.font.size = Pt(22)
    right_para.font.bold = True
    right_para.font.color.rgb = PRIMARY_BLUE

    # Right content
    right_content = slide.shapes.add_textbox(Inches(5), Inches(2.1), Inches(4.5), Inches(4))
    right_content_frame = right_content.text_frame
    right_content_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            para = right_content_frame.paragraphs[0]
        else:
            para = right_content_frame.add_paragraph()
        para.text = f"• {item}"
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY
        para.space_after = Pt(8)

def add_pricing_slide(prs, lang_data):
    """Add pricing slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_BLUE
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = lang_data["pricing_title"]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Three pricing boxes
    plans = lang_data["pricing_plans"]
    box_width = Inches(3)
    box_height = Inches(4.5)
    start_x = Inches(0.4)
    y = Inches(1.5)
    gap = Inches(0.2)

    for i, plan in enumerate(plans):
        x = start_x + (box_width + gap) * i

        # Box background
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
        if i == 1:  # Highlight middle plan
            box.fill.solid()
            box.fill.fore_color.rgb = PRIMARY_BLUE
            text_color = WHITE
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT_GRAY
            text_color = DARK_GRAY
        box.line.fill.background()

        # Plan name
        name_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), box_width - Inches(0.4), Inches(0.5))
        name_frame = name_box.text_frame
        name_para = name_frame.paragraphs[0]
        name_para.text = plan["name"]
        name_para.font.size = Pt(20)
        name_para.font.bold = True
        name_para.font.color.rgb = text_color
        name_para.alignment = PP_ALIGN.CENTER

        # Price
        price_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), box_width - Inches(0.4), Inches(0.5))
        price_frame = price_box.text_frame
        price_para = price_frame.paragraphs[0]
        price_para.text = plan["price"]
        price_para.font.size = Pt(24)
        price_para.font.bold = True
        price_para.font.color.rgb = text_color
        price_para.alignment = PP_ALIGN.CENTER

        # Features
        features_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.3), box_width - Inches(0.3), Inches(3))
        features_frame = features_box.text_frame
        features_frame.word_wrap = True
        for j, feature in enumerate(plan["features"]):
            if j == 0:
                para = features_frame.paragraphs[0]
            else:
                para = features_frame.add_paragraph()
            para.text = f"✓ {feature}"
            para.font.size = Pt(12)
            para.font.color.rgb = text_color
            para.space_after = Pt(4)

def add_partner_slide(prs, lang_data):
    """Add partner benefits slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_BLUE
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = lang_data["partner_title"]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Partner tiers table
    tiers = lang_data["partner_tiers"]

    # Commission highlight
    commission_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
    comm_frame = commission_box.text_frame
    comm_para = comm_frame.paragraphs[0]
    comm_para.text = lang_data["commission_highlight"]
    comm_para.font.size = Pt(24)
    comm_para.font.bold = True
    comm_para.font.color.rgb = PRIMARY_BLUE
    comm_para.alignment = PP_ALIGN.CENTER

    # Tier boxes
    box_width = Inches(3)
    start_x = Inches(0.4)
    y = Inches(2.5)
    gap = Inches(0.2)

    colors = [LIGHT_GRAY, RGBColor(253, 224, 71), RGBColor(229, 231, 235)]  # Silver, Gold, Platinum

    for i, tier in enumerate(tiers):
        x = start_x + (box_width + gap) * i

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.fill.background()

        # Tier name
        name_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), box_width - Inches(0.4), Inches(0.5))
        name_frame = name_box.text_frame
        name_para = name_frame.paragraphs[0]
        name_para.text = tier["name"]
        name_para.font.size = Pt(20)
        name_para.font.bold = True
        name_para.font.color.rgb = DARK_GRAY
        name_para.alignment = PP_ALIGN.CENTER

        # Revenue
        rev_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), box_width - Inches(0.4), Inches(0.4))
        rev_frame = rev_box.text_frame
        rev_para = rev_frame.paragraphs[0]
        rev_para.text = tier["revenue"]
        rev_para.font.size = Pt(14)
        rev_para.font.color.rgb = DARK_GRAY
        rev_para.alignment = PP_ALIGN.CENTER

        # Commission
        comm_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.2), box_width - Inches(0.4), Inches(0.5))
        comm_frame = comm_box.text_frame
        comm_para = comm_frame.paragraphs[0]
        comm_para.text = tier["commission"]
        comm_para.font.size = Pt(28)
        comm_para.font.bold = True
        comm_para.font.color.rgb = PRIMARY_BLUE
        comm_para.alignment = PP_ALIGN.CENTER

        # Benefits
        benefits_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.8), box_width - Inches(0.3), Inches(1.5))
        benefits_frame = benefits_box.text_frame
        benefits_frame.word_wrap = True
        for j, benefit in enumerate(tier["benefits"]):
            if j == 0:
                para = benefits_frame.paragraphs[0]
            else:
                para = benefits_frame.add_paragraph()
            para.text = f"✓ {benefit}"
            para.font.size = Pt(11)
            para.font.color.rgb = DARK_GRAY
            para.space_after = Pt(2)

def add_contact_slide(prs, lang_data):
    """Add contact/thank you slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_BLUE
    background.line.fill.background()

    # Thank you text
    thank_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    thank_frame = thank_box.text_frame
    thank_para = thank_frame.paragraphs[0]
    thank_para.text = lang_data["thank_you"]
    thank_para.font.size = Pt(44)
    thank_para.font.bold = True
    thank_para.font.color.rgb = WHITE
    thank_para.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(2))
    contact_frame = contact_box.text_frame

    contact_para = contact_frame.paragraphs[0]
    contact_para.text = "sales@datamigrate.ai"
    contact_para.font.size = Pt(24)
    contact_para.font.color.rgb = WHITE
    contact_para.alignment = PP_ALIGN.CENTER

    web_para = contact_frame.add_paragraph()
    web_para.text = "www.datamigrate.ai"
    web_para.font.size = Pt(24)
    web_para.font.color.rgb = WHITE
    web_para.alignment = PP_ALIGN.CENTER

    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tag_frame = tag_box.text_frame
    tag_para = tag_frame.paragraphs[0]
    tag_para.text = lang_data["tagline"]
    tag_para.font.size = Pt(16)
    tag_para.font.italic = True
    tag_para.font.color.rgb = WHITE
    tag_para.alignment = PP_ALIGN.CENTER


def create_presentation(lang_code, lang_data, output_dir):
    """Create a complete presentation for the given language"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "DataMigrate AI", lang_data["subtitle"], lang_data)

    # Slide 2: The Problem
    add_content_slide(prs, lang_data["problem_title"], lang_data["problem_bullets"])

    # Slide 3: The Solution
    add_content_slide(prs, lang_data["solution_title"], lang_data["solution_bullets"])

    # Slide 4: How It Works
    add_content_slide(prs, lang_data["how_title"], lang_data["how_bullets"])

    # Slide 5: AI Agents
    add_two_column_slide(
        prs,
        lang_data["agents_title"],
        lang_data["agents_left_title"],
        lang_data["agents_left"],
        lang_data["agents_right_title"],
        lang_data["agents_right"]
    )

    # Slide 6: Warehouses & Security
    add_two_column_slide(
        prs,
        lang_data["warehouse_security_title"],
        lang_data["warehouse_title"],
        lang_data["warehouses"],
        lang_data["security_title"],
        lang_data["security_items"]
    )

    # Slide 7: Benefits
    add_content_slide(prs, lang_data["benefits_title"], lang_data["benefits_bullets"])

    # Slide 8: Pricing
    add_pricing_slide(prs, lang_data)

    # Slide 9: Partner Benefits
    add_partner_slide(prs, lang_data)

    # Slide 10: Contact
    add_contact_slide(prs, lang_data)

    # Save
    output_path = os.path.join(output_dir, lang_code, f"DataMigrate-AI-Sales-{lang_code.upper()}.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path


# Language data
LANGUAGES = {
    "en": {
        "subtitle": "Transform Your Legacy Data Infrastructure with AI-Powered Migration",
        "problem_title": "The Problem",
        "problem_bullets": [
            "Manual SQL Server migrations take weeks or months",
            "Human errors in data transformation cause costly mistakes",
            "Requires specialized DBAs and data engineers",
            "Legacy systems often lack proper documentation",
            "Manual testing misses edge cases and data quality issues"
        ],
        "solution_title": "The Solution: DataMigrate AI",
        "solution_bullets": [
            "First AI-powered MSSQL to dbt migration platform",
            "Automated schema analysis and transformation",
            "AI-generated dbt models with best practices",
            "Comprehensive testing and documentation",
            "Support for all major cloud data warehouses"
        ],
        "how_title": "How It Works",
        "how_bullets": [
            "1. CONNECT - Secure connection to your SQL Server database",
            "2. ANALYZE - AI agents analyze schema, relationships, and patterns",
            "3. TRANSFORM - Generate optimized dbt models and tests",
            "4. DEPLOY - Export to Snowflake, BigQuery, Fabric, Databricks, or Redshift"
        ],
        "agents_title": "AI Agents",
        "agents_left_title": "Data Preparation",
        "agents_left": ["DataPrep Agent - Profiling & cleansing", "ML Fine-Tuning - Pattern recognition", "Data Quality - Automated testing"],
        "agents_right_title": "Documentation & Analytics",
        "agents_right": ["Documentation Agent - Auto-generated docs", "BI Agent - Analytics-ready models", "Lineage tracking & compliance"],
        "warehouse_security_title": "Deploy Anywhere, Securely",
        "warehouse_title": "Target Warehouses",
        "warehouses": ["Snowflake", "Google BigQuery", "Microsoft Fabric", "Databricks", "Amazon Redshift"],
        "security_title": "Enterprise Security",
        "security_items": ["TLS 1.3 + AES-256 encryption", "JWT authentication & RBAC", "GDPR compliant", "Audit logging", "On-premises available"],
        "benefits_title": "Customer Benefits",
        "benefits_bullets": [
            "80% Faster Migrations - Days instead of months",
            "95% Fewer Errors - Automated testing catches issues",
            "60% Cost Reduction - Less manual work",
            "100% Documentation - Everything fully documented",
            "Zero Lock-In - Standard dbt output"
        ],
        "pricing_title": "Pricing Plans",
        "pricing_plans": [
            {"name": "Starter", "price": "$499/mo", "features": ["Up to 50 tables", "3 migrations/month", "Email support", "Basic warehouses"]},
            {"name": "Professional", "price": "$1,499/mo", "features": ["Up to 200 tables", "Unlimited migrations", "Priority support", "All warehouses", "All AI agents"]},
            {"name": "Enterprise", "price": "Custom", "features": ["Unlimited tables", "Dedicated support", "On-premises option", "SLA guarantees", "Custom integrations"]}
        ],
        "partner_title": "Partner Benefits",
        "commission_highlight": "15-20% Recurring Commission on All Deals",
        "partner_tiers": [
            {"name": "Silver", "revenue": "$0 - $50K", "commission": "15%", "benefits": ["Sales materials", "Product training", "Demo access"]},
            {"name": "Gold", "revenue": "$50K - $200K", "commission": "18%", "benefits": ["Priority support", "Lead sharing", "Co-marketing"]},
            {"name": "Platinum", "revenue": "$200K+", "commission": "20%", "benefits": ["Dedicated manager", "Joint events", "Custom solutions"]}
        ],
        "thank_you": "Let's Transform Your Data",
        "tagline": "AI-Powered | Secure | Multilingual | Enterprise-Ready"
    },
    "es": {
        "subtitle": "Transforma tu Infraestructura de Datos con Migración Impulsada por IA",
        "problem_title": "El Problema",
        "problem_bullets": [
            "Las migraciones manuales de SQL Server toman semanas o meses",
            "Los errores humanos en transformación de datos causan costosos errores",
            "Requiere DBAs y ingenieros de datos especializados",
            "Los sistemas heredados carecen de documentación adecuada",
            "Las pruebas manuales pierden casos extremos y problemas de calidad"
        ],
        "solution_title": "La Solución: DataMigrate AI",
        "solution_bullets": [
            "Primera plataforma de migración MSSQL a dbt impulsada por IA",
            "Análisis y transformación automática de esquemas",
            "Modelos dbt generados por IA con mejores prácticas",
            "Pruebas y documentación completas",
            "Soporte para todos los principales data warehouses en la nube"
        ],
        "how_title": "Cómo Funciona",
        "how_bullets": [
            "1. CONECTAR - Conexión segura a tu base de datos SQL Server",
            "2. ANALIZAR - Los agentes IA analizan esquemas, relaciones y patrones",
            "3. TRANSFORMAR - Generar modelos dbt optimizados y pruebas",
            "4. DESPLEGAR - Exportar a Snowflake, BigQuery, Fabric, Databricks o Redshift"
        ],
        "agents_title": "Agentes de IA",
        "agents_left_title": "Preparación de Datos",
        "agents_left": ["Agente DataPrep - Perfilado y limpieza", "ML Fine-Tuning - Reconocimiento de patrones", "Calidad de Datos - Pruebas automatizadas"],
        "agents_right_title": "Documentación y Analítica",
        "agents_right": ["Agente de Documentación - Docs auto-generados", "Agente BI - Modelos listos para analítica", "Seguimiento de linaje y cumplimiento"],
        "warehouse_security_title": "Despliega en Cualquier Lugar, de Forma Segura",
        "warehouse_title": "Warehouses Destino",
        "warehouses": ["Snowflake", "Google BigQuery", "Microsoft Fabric", "Databricks", "Amazon Redshift"],
        "security_title": "Seguridad Empresarial",
        "security_items": ["Cifrado TLS 1.3 + AES-256", "Autenticación JWT y RBAC", "Cumple con GDPR", "Registro de auditoría", "Disponible on-premises"],
        "benefits_title": "Beneficios para el Cliente",
        "benefits_bullets": [
            "80% Más Rápido - Días en lugar de meses",
            "95% Menos Errores - Pruebas automatizadas detectan problemas",
            "60% Reducción de Costos - Menos trabajo manual",
            "100% Documentación - Todo completamente documentado",
            "Sin Lock-In - Salida estándar dbt"
        ],
        "pricing_title": "Planes de Precios",
        "pricing_plans": [
            {"name": "Starter", "price": "$499/mes", "features": ["Hasta 50 tablas", "3 migraciones/mes", "Soporte por email", "Warehouses básicos"]},
            {"name": "Professional", "price": "$1,499/mes", "features": ["Hasta 200 tablas", "Migraciones ilimitadas", "Soporte prioritario", "Todos los warehouses", "Todos los agentes IA"]},
            {"name": "Enterprise", "price": "Personalizado", "features": ["Tablas ilimitadas", "Soporte dedicado", "Opción on-premises", "Garantías SLA", "Integraciones personalizadas"]}
        ],
        "partner_title": "Beneficios para Socios",
        "commission_highlight": "15-20% Comisión Recurrente en Todos los Negocios",
        "partner_tiers": [
            {"name": "Silver", "revenue": "$0 - $50K", "commission": "15%", "benefits": ["Materiales de ventas", "Capacitación", "Acceso a demos"]},
            {"name": "Gold", "revenue": "$50K - $200K", "commission": "18%", "benefits": ["Soporte prioritario", "Leads compartidos", "Co-marketing"]},
            {"name": "Platinum", "revenue": "$200K+", "commission": "20%", "benefits": ["Gerente dedicado", "Eventos conjuntos", "Soluciones custom"]}
        ],
        "thank_you": "Transformemos tus Datos",
        "tagline": "Impulsado por IA | Seguro | Multilingüe | Listo para Empresas"
    },
    "pt": {
        "subtitle": "Transforme sua Infraestrutura de Dados com Migração Impulsionada por IA",
        "problem_title": "O Problema",
        "problem_bullets": [
            "Migrações manuais de SQL Server levam semanas ou meses",
            "Erros humanos na transformação de dados causam erros custosos",
            "Requer DBAs e engenheiros de dados especializados",
            "Sistemas legados frequentemente não têm documentação adequada",
            "Testes manuais perdem casos extremos e problemas de qualidade"
        ],
        "solution_title": "A Solução: DataMigrate AI",
        "solution_bullets": [
            "Primeira plataforma de migração MSSQL para dbt com IA",
            "Análise e transformação automática de esquemas",
            "Modelos dbt gerados por IA com melhores práticas",
            "Testes e documentação abrangentes",
            "Suporte para todos os principais data warehouses na nuvem"
        ],
        "how_title": "Como Funciona",
        "how_bullets": [
            "1. CONECTAR - Conexão segura ao seu banco SQL Server",
            "2. ANALISAR - Agentes IA analisam esquemas, relacionamentos e padrões",
            "3. TRANSFORMAR - Gerar modelos dbt otimizados e testes",
            "4. IMPLANTAR - Exportar para Snowflake, BigQuery, Fabric, Databricks ou Redshift"
        ],
        "agents_title": "Agentes de IA",
        "agents_left_title": "Preparação de Dados",
        "agents_left": ["Agente DataPrep - Perfilamento e limpeza", "ML Fine-Tuning - Reconhecimento de padrões", "Qualidade de Dados - Testes automatizados"],
        "agents_right_title": "Documentação e Analytics",
        "agents_right": ["Agente de Documentação - Docs auto-gerados", "Agente BI - Modelos prontos para analytics", "Rastreamento de linhagem e conformidade"],
        "warehouse_security_title": "Implante em Qualquer Lugar, com Segurança",
        "warehouse_title": "Warehouses Destino",
        "warehouses": ["Snowflake", "Google BigQuery", "Microsoft Fabric", "Databricks", "Amazon Redshift"],
        "security_title": "Segurança Empresarial",
        "security_items": ["Criptografia TLS 1.3 + AES-256", "Autenticação JWT e RBAC", "Conformidade com GDPR", "Log de auditoria", "Disponível on-premises"],
        "benefits_title": "Benefícios para o Cliente",
        "benefits_bullets": [
            "80% Mais Rápido - Dias em vez de meses",
            "95% Menos Erros - Testes automatizados detectam problemas",
            "60% Redução de Custos - Menos trabalho manual",
            "100% Documentação - Tudo completamente documentado",
            "Sem Lock-In - Saída padrão dbt"
        ],
        "pricing_title": "Planos de Preços",
        "pricing_plans": [
            {"name": "Starter", "price": "$499/mês", "features": ["Até 50 tabelas", "3 migrações/mês", "Suporte por email", "Warehouses básicos"]},
            {"name": "Professional", "price": "$1.499/mês", "features": ["Até 200 tabelas", "Migrações ilimitadas", "Suporte prioritário", "Todos os warehouses", "Todos os agentes IA"]},
            {"name": "Enterprise", "price": "Personalizado", "features": ["Tabelas ilimitadas", "Suporte dedicado", "Opção on-premises", "Garantias SLA", "Integrações personalizadas"]}
        ],
        "partner_title": "Benefícios para Parceiros",
        "commission_highlight": "15-20% Comissão Recorrente em Todos os Negócios",
        "partner_tiers": [
            {"name": "Silver", "revenue": "$0 - $50K", "commission": "15%", "benefits": ["Materiais de vendas", "Treinamento", "Acesso a demos"]},
            {"name": "Gold", "revenue": "$50K - $200K", "commission": "18%", "benefits": ["Suporte prioritário", "Leads compartilhados", "Co-marketing"]},
            {"name": "Platinum", "revenue": "$200K+", "commission": "20%", "benefits": ["Gerente dedicado", "Eventos conjuntos", "Soluções custom"]}
        ],
        "thank_you": "Vamos Transformar seus Dados",
        "tagline": "Impulsionado por IA | Seguro | Multilíngue | Pronto para Empresas"
    },
    "da": {
        "subtitle": "Transformer din Datainfrastruktur med AI-drevet Migration",
        "problem_title": "Problemet",
        "problem_bullets": [
            "Manuelle SQL Server-migreringer tager uger eller måneder",
            "Menneskelige fejl i datatransformation forårsager dyre fejl",
            "Kræver specialiserede DBA'er og dataingeniører",
            "Ældre systemer mangler ofte ordentlig dokumentation",
            "Manuel test overser edge cases og datakvalitetsproblemer"
        ],
        "solution_title": "Løsningen: DataMigrate AI",
        "solution_bullets": [
            "Første AI-drevne MSSQL til dbt migrationsplatform",
            "Automatiseret skemaanalyse og transformation",
            "AI-genererede dbt-modeller med best practices",
            "Omfattende test og dokumentation",
            "Support til alle større cloud data warehouses"
        ],
        "how_title": "Sådan Fungerer Det",
        "how_bullets": [
            "1. FORBIND - Sikker forbindelse til din SQL Server database",
            "2. ANALYSER - AI-agenter analyserer skemaer, relationer og mønstre",
            "3. TRANSFORMER - Generer optimerede dbt-modeller og tests",
            "4. DEPLOY - Eksporter til Snowflake, BigQuery, Fabric, Databricks eller Redshift"
        ],
        "agents_title": "AI-Agenter",
        "agents_left_title": "Dataforberedelse",
        "agents_left": ["DataPrep Agent - Profilering og rensning", "ML Fine-Tuning - Mønstergenkendelse", "Datakvalitet - Automatiseret test"],
        "agents_right_title": "Dokumentation og Analytics",
        "agents_right": ["Dokumentationsagent - Auto-genererede docs", "BI Agent - Analytics-klare modeller", "Lineage tracking og compliance"],
        "warehouse_security_title": "Deploy Overalt, Sikkert",
        "warehouse_title": "Mål Warehouses",
        "warehouses": ["Snowflake", "Google BigQuery", "Microsoft Fabric", "Databricks", "Amazon Redshift"],
        "security_title": "Virksomhedssikkerhed",
        "security_items": ["TLS 1.3 + AES-256 kryptering", "JWT-godkendelse og RBAC", "GDPR-kompatibel", "Revisionslogning", "On-premises tilgængelig"],
        "benefits_title": "Kundefordele",
        "benefits_bullets": [
            "80% Hurtigere - Dage i stedet for måneder",
            "95% Færre Fejl - Automatiseret test fanger problemer",
            "60% Omkostningsreduktion - Mindre manuelt arbejde",
            "100% Dokumentation - Alt fuldt dokumenteret",
            "Ingen Lock-In - Standard dbt output"
        ],
        "pricing_title": "Prisplaner",
        "pricing_plans": [
            {"name": "Starter", "price": "$499/md", "features": ["Op til 50 tabeller", "3 migreringer/måned", "Email support", "Basis warehouses"]},
            {"name": "Professional", "price": "$1.499/md", "features": ["Op til 200 tabeller", "Ubegrænsede migreringer", "Prioriteret support", "Alle warehouses", "Alle AI-agenter"]},
            {"name": "Enterprise", "price": "Tilpasset", "features": ["Ubegrænsede tabeller", "Dedikeret support", "On-premises mulighed", "SLA-garantier", "Tilpassede integrationer"]}
        ],
        "partner_title": "Partnerfordele",
        "commission_highlight": "15-20% Tilbagevendende Provision på Alle Handler",
        "partner_tiers": [
            {"name": "Silver", "revenue": "$0 - $50K", "commission": "15%", "benefits": ["Salgsmaterialer", "Produkttræning", "Demo-adgang"]},
            {"name": "Gold", "revenue": "$50K - $200K", "commission": "18%", "benefits": ["Prioriteret support", "Lead-deling", "Co-marketing"]},
            {"name": "Platinum", "revenue": "$200K+", "commission": "20%", "benefits": ["Dedikeret manager", "Fælles events", "Custom løsninger"]}
        ],
        "thank_you": "Lad os Transformere dine Data",
        "tagline": "AI-drevet | Sikker | Flersproget | Virksomhedsklar"
    },
    "de": {
        "subtitle": "Transformieren Sie Ihre Dateninfrastruktur mit KI-gestützter Migration",
        "problem_title": "Das Problem",
        "problem_bullets": [
            "Manuelle SQL Server-Migrationen dauern Wochen oder Monate",
            "Menschliche Fehler bei der Datentransformation verursachen kostspielige Fehler",
            "Erfordert spezialisierte DBAs und Dateningenieure",
            "Legacy-Systeme haben oft keine ordentliche Dokumentation",
            "Manuelle Tests übersehen Grenzfälle und Datenqualitätsprobleme"
        ],
        "solution_title": "Die Lösung: DataMigrate AI",
        "solution_bullets": [
            "Erste KI-gestützte MSSQL zu dbt Migrationsplattform",
            "Automatisierte Schemaanalyse und Transformation",
            "KI-generierte dbt-Modelle mit Best Practices",
            "Umfassende Tests und Dokumentation",
            "Unterstützung für alle großen Cloud Data Warehouses"
        ],
        "how_title": "So Funktioniert Es",
        "how_bullets": [
            "1. VERBINDEN - Sichere Verbindung zu Ihrer SQL Server Datenbank",
            "2. ANALYSIEREN - KI-Agenten analysieren Schemas, Beziehungen und Muster",
            "3. TRANSFORMIEREN - Optimierte dbt-Modelle und Tests generieren",
            "4. BEREITSTELLEN - Export nach Snowflake, BigQuery, Fabric, Databricks oder Redshift"
        ],
        "agents_title": "KI-Agenten",
        "agents_left_title": "Datenvorbereitung",
        "agents_left": ["DataPrep Agent - Profiling und Bereinigung", "ML Fine-Tuning - Mustererkennung", "Datenqualität - Automatisierte Tests"],
        "agents_right_title": "Dokumentation und Analytics",
        "agents_right": ["Dokumentationsagent - Auto-generierte Docs", "BI Agent - Analytics-bereite Modelle", "Lineage-Tracking und Compliance"],
        "warehouse_security_title": "Überall Bereitstellen, Sicher",
        "warehouse_title": "Ziel-Warehouses",
        "warehouses": ["Snowflake", "Google BigQuery", "Microsoft Fabric", "Databricks", "Amazon Redshift"],
        "security_title": "Unternehmenssicherheit",
        "security_items": ["TLS 1.3 + AES-256 Verschlüsselung", "JWT-Authentifizierung und RBAC", "DSGVO-konform", "Audit-Protokollierung", "On-Premises verfügbar"],
        "benefits_title": "Kundenvorteile",
        "benefits_bullets": [
            "80% Schneller - Tage statt Monate",
            "95% Weniger Fehler - Automatisierte Tests erkennen Probleme",
            "60% Kostenreduzierung - Weniger manuelle Arbeit",
            "100% Dokumentation - Alles vollständig dokumentiert",
            "Kein Lock-In - Standard dbt Output"
        ],
        "pricing_title": "Preispläne",
        "pricing_plans": [
            {"name": "Starter", "price": "€499/Mo", "features": ["Bis zu 50 Tabellen", "3 Migrationen/Monat", "E-Mail-Support", "Basis Warehouses"]},
            {"name": "Professional", "price": "€1.499/Mo", "features": ["Bis zu 200 Tabellen", "Unbegrenzte Migrationen", "Prioritäts-Support", "Alle Warehouses", "Alle KI-Agenten"]},
            {"name": "Enterprise", "price": "Individuell", "features": ["Unbegrenzte Tabellen", "Dedizierter Support", "On-Premises Option", "SLA-Garantien", "Individuelle Integrationen"]}
        ],
        "partner_title": "Partnervorteile",
        "commission_highlight": "15-20% Wiederkehrende Provision auf Alle Geschäfte",
        "partner_tiers": [
            {"name": "Silver", "revenue": "€0 - €50K", "commission": "15%", "benefits": ["Vertriebsmaterialien", "Produktschulung", "Demo-Zugang"]},
            {"name": "Gold", "revenue": "€50K - €200K", "commission": "18%", "benefits": ["Prioritäts-Support", "Lead-Sharing", "Co-Marketing"]},
            {"name": "Platinum", "revenue": "€200K+", "commission": "20%", "benefits": ["Dedizierter Manager", "Gemeinsame Events", "Custom Lösungen"]}
        ],
        "thank_you": "Transformieren wir Ihre Daten",
        "tagline": "KI-gestützt | Sicher | Mehrsprachig | Unternehmensfähig"
    }
}


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))

    print("DataMigrate AI - Sales Presentation Generator")
    print("=" * 50)

    for lang_code, lang_data in LANGUAGES.items():
        try:
            create_presentation(lang_code, lang_data, script_dir)
        except Exception as e:
            print(f"Error creating {lang_code}: {e}")

    print("\nAll presentations created successfully!")
    print(f"Output directory: {script_dir}")


if __name__ == "__main__":
    main()
